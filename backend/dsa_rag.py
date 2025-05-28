"""
RAG Module
This module implements the Retrieval-Augmented Generation (RAG) system.
It handles document indexing, text extraction from various file formats,
and integration with the vector database for semantic search.
"""

# Import docx module for Word file manipulation
import docx

# Import PyPDF2 module for PDF file manipulation
import PyPDF2

# Import Presentation from pptx package for PowerPoint file manipulation
from pptx import Presentation

# Import TokenTextSplitter for text tokenization
from langchain_text_splitters import TokenTextSplitter

# Import HuggingFaceEmbeddings for creating embeddings
from langchain_huggingface import HuggingFaceEmbeddings

# Import Qdrant classes for vector database operations
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# Import Qdrant for vector database integration
from langchain_qdrant import Qdrant

# Import storage module
from storage import get_storage

# Import required modules
import os


def dsa_carrega_texto_word(arquivoname):
    """
    Load text content from a Word document.
    
    Args:
        arquivoname (str): Path to the Word document
        
    Returns:
        str: Extracted text content from the document
    """
    doc = docx.Document(arquivoname)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)


def dsa_carrega_texto_pptx(arquivoname):
    """
    Load text content from a PowerPoint presentation.
    
    Args:
        arquivoname (str): Path to the PowerPoint file
        
    Returns:
        str: Extracted text content from the presentation
    """
    prs = Presentation(arquivoname)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)


def main_indexing(storage_config):
    """
    Main function for document indexing.
    Processes documents from storage, extracts text, and indexes them in the vector database.
    
    Args:
        storage_config (dict): Configuration for document storage
            - storage_type: Type of storage ('local', 's3', or 'http')
            - base_path: Base path for local storage
            - bucket_name: S3 bucket name (for S3 storage)
            - region_name: AWS region (for S3 storage)
            - endpoint_url: S3 endpoint URL (for S3 storage)
            - base_url: Base URL for HTTP storage
    """
    # Initialize storage
    storage = get_storage(**storage_config)
    
    # Define model parameters for embeddings
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
    model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}

    # Initialize HuggingFace embeddings
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )

    # Initialize Qdrant client
    client = QdrantClient("http://qdrant:6333")
    collection_name = "DSAVectorDB"

    # Delete collection if it exists
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    # Create new collection with specified parameters
    client.create_collection(
        collection_name,
        vectors_config=VectorParams(size=768, distance=Distance.DOT)
    )

    # Initialize Qdrant instance
    qdrant = Qdrant(client, collection_name, hf)

    print("\nIndexing documents...\n")

    # Get list of all documents
    lista_arquivos = storage.list_documents()
    
    # Process each file in the list
    for arquivo in lista_arquivos:
        
        try:
            # Get document from storage
            temp_file = storage.get_document(arquivo)
            
            try:
                arquivo_content = ""
                
                # Process PDF files
                if arquivo.endswith(".pdf"):
                    print("Indexing: " + arquivo)
                    reader = PyPDF2.PdfReader(temp_file)
                    for page in reader.pages:
                        arquivo_content += " " + page.extract_text()
                
                # Process text files
                elif arquivo.endswith(".txt"):
                    print("Indexing: " + arquivo)
                    with open(temp_file, 'r') as f:
                        arquivo_content = f.read()
                
                # Process Word documents
                elif arquivo.endswith(".docx"):
                    print("Indexing: " + arquivo)
                    arquivo_content = dsa_carrega_texto_word(temp_file)
                
                # Process PowerPoint presentations
                elif arquivo.endswith(".pptx"):
                    print("Indexing: " + arquivo)
                    arquivo_content = dsa_carrega_texto_pptx(temp_file)
                
                else:
                    continue

                # Initialize text splitter
                text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
                textos = text_splitter.split_text(arquivo_content)
                metadata = [{"path": arquivo} for _ in textos]
                qdrant.add_texts(textos, metadatas=metadata)

            finally:

                # Remove temporary file
                os.unlink(temp_file)

        except Exception as e:
            print(f"Process failed for file {arquivo}: {e}")

    print("\nIndexing Completed!\n")


if __name__ == "__main__":

    # Configure storage
    storage_config = {
        'storage_type': os.getenv('STORAGE_TYPE', 'local'),
        'base_path': os.getenv('DOCUMENTS_PATH', '/app/documents')
    }
    
    # Add S3 specific configurations if needed
    if storage_config['storage_type'] == 's3':
        storage_config.update({
            'bucket_name': os.getenv('S3_BUCKET_NAME'),
            'region_name': os.getenv('AWS_REGION'),
            'endpoint_url': os.getenv('S3_ENDPOINT_URL')
        })
    # Add HTTP base URL if needed
    elif storage_config['storage_type'] == 'http':
        storage_config['base_url'] = os.getenv('DOCUMENT_STORAGE_URL', 'http://document_storage:8080')
    
    main_indexing(storage_config)




