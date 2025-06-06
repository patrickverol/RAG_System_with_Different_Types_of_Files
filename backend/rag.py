"""
RAG Module
This module implements the Retrieval-Augmented Generation (RAG) system.
It handles document indexing, text extraction from various file formats,
and integration with the vector database for semantic search.
"""

#import csv

import csv



# Import docx module for Word file manipulation
import docx

# Import PyPDF2 module for PDF file manipulation
import PyPDF2

# Import Presentation from pptx package for PowerPoint file manipulation
from pptx import Presentation

# Import TokenTextSplitter for text tokenization
from langchain_text_splitters import TokenTextSplitter

# Import HuggingFaceEmbeddings for creating embeddings
from langchain_huggingface import HuggingFaceEmbeddings

# Import Qdrant classes for vector database operations
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# Import QdrantVectorStore for vector database integration
from langchain_qdrant import QdrantVectorStore

# Import storage module
from storage import get_storage

# Import required modules
import os
import tempfile


def carrega_texto_word(arquivoname):
    """
    Load text content from a Word document.
    
    Args:
        arquivoname (str): Path to the Word document
        
    Returns:
        str: Extracted text content from the document
    """
    doc = docx.Document(arquivoname)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)


def carrega_texto_pptx(arquivoname):
    """
    Load text content from a PowerPoint presentation.
    
    Args:
        arquivoname (str): Path to the PowerPoint file
        
    Returns:
        str: Extracted text content from the presentation
    """
    prs = Presentation(arquivoname)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)


def main_indexing(storage_config):
    """
    Main function for document indexing.
    Processes documents from storage, extracts text, and indexes them in the vector database.
    
    Args:
        storage_config (dict): Configuration for document storage
            - storage_type: Type of storage ('local', 's3', or 'http')
            - base_path: Base path for local storage
            - bucket_name: S3 bucket name (for S3 storage)
            - region_name: AWS region (for S3 storage)
            - endpoint_url: S3 endpoint URL (for S3 storage)
            - base_url: Base URL for HTTP storage
    """
    print(f"\nStorage config: {storage_config}")
    
    # Initialize storage
    storage = get_storage(**storage_config)
    
    # Define model parameters for embeddings
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
    model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}

    # Initialize HuggingFace embeddings
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )

    # Initialize Qdrant client
    client = QdrantClient("http://qdrant:6333")
    collection_name = "RAGVectorDB"

    # Delete collection if it exists
    if client.collection_exists(collection_name):
        print(f"Deleting existing collection: {collection_name}")
        client.delete_collection(collection_name)

    # Create new collection with specified parameters
    print(f"Creating new collection: {collection_name}")
    client.create_collection(
        collection_name,
        vectors_config=VectorParams(size=768, distance=Distance.COSINE)
    )

    # Initialize Qdrant instance
    qdrant = QdrantVectorStore(
        client=client,
        collection_name=collection_name,
        embedding=hf,
        distance=Distance.COSINE
    )

    print("\nIndexing documents...\n")

    # Get list of all documents
    lista_arquivos = storage.list_documents()
    print(f"Found {len(lista_arquivos)} documents to index")
    
    # Process each file in the list
    for arquivo in lista_arquivos:
        print(f"\nProcessing file: {arquivo}")
        
        try:
            # Get document from storage
            temp_file = storage.get_document(arquivo)
            print(f"Retrieved file from storage: {temp_file}")
            
            try:
                arquivo_content = ""
                
                # Process PDF files
                if arquivo.endswith(".pdf"):
                    print("Processing PDF file")
                    reader = PyPDF2.PdfReader(temp_file)
                    for page in reader.pages:
                        arquivo_content += " " + page.extract_text()
                
                # Process text files
                elif arquivo.endswith(".txt"):
                    print("Processing text file")
                    with open(temp_file, 'r') as f:
                        arquivo_content = f.read()
                
                # Process Word documents
                elif arquivo.endswith(".docx"):
                    print("Processing Word document")
                    arquivo_content = carrega_texto_word(temp_file)
                
                # Process PowerPoint presentations
                elif arquivo.endswith(".pptx"):
                    print("Processing PowerPoint presentation")
                    arquivo_content = carrega_texto_pptx(temp_file)

                # Process CSV files
                elif arquivo.endswith(".csv"):
                    print("Processing CSV file ....")
                    rows = []
                    with open(temp_file, 'r', encoding='utf-8') as csvfile:
                        reader = csv.reader(csvfile)
                        for row in reader:
                            # Join each cell with a space, then append
                            rows.append(" ".join(row))
                    arquivo_content = "\n".join(rows)
                    print(f"Extracted {len(rows)} CSV rows as text")
                
                else:
                    print(f"Skipping unsupported file type: {arquivo}")
                    continue

                print(f"Extracted content length: {len(arquivo_content)} characters")

                # Initialize text splitter
                text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
                textos = text_splitter.split_text(arquivo_content)
                print(f"Split into {len(textos)} chunks")
                
                metadata = [{"path": arquivo} for _ in textos]
                qdrant.add_texts(textos, metadatas=metadata)
                print(f"Successfully indexed {len(textos)} chunks")

            finally:
                # Only delete the temporary file if it's in the temp directory
                if tempfile.gettempdir() in temp_file:
                    os.unlink(temp_file)
                    print(f"Cleaned up temporary file: {temp_file}")

        except Exception as e:
            print(f"Process failed for file {arquivo}: {str(e)}")

    print("\nIndexing Completed!\n")


if __name__ == "__main__":

    # Get storage type from environment variable
    storage_type = os.getenv('STORAGE_TYPE', 'http')
    
    # Configure storage based on type
    if storage_type == 's3':
        storage_config = {
            'storage_type': 's3',
            'bucket_name': os.getenv('S3_BUCKET_NAME'),
            'region_name': os.getenv('AWS_REGION'),
            'endpoint_url': os.getenv('S3_ENDPOINT_URL')
        }
    elif storage_type == 'local':
        storage_config = {
            'storage_type': 'local',
            'base_path': os.getenv('DOCUMENTS_PATH', '/app/documents')
        }
    else:  # http
        storage_config = {
            'storage_type': 'http',
            'base_url': os.getenv('DOCUMENT_STORAGE_URL', 'http://document_storage:8080')
        }
    
    main_indexing(storage_config)